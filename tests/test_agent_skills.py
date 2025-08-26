"""Adapted tests for simple_docker_runtime agent_skills.

Source: OpenHands tests/unit/test_agent_skill.py
- Imports updated to `simple_docker_runtime.plugins.agent_skills...`
- Fixture resets CURRENT_FILE in the file_ops module (not in agentskills)
- PDF test is skipped if `reportlab` is not available
"""

from __future__ import annotations

import contextlib
import io
import sys
from pathlib import Path

import pytest

from simple_openhands.plugins.agent_skills.file_ops.file_ops import (
    WINDOW,
    _print_window,
    find_file,
    goto_line,
    open_file,
    scroll_down,
    scroll_up,
    search_dir,
    search_file,
)
from simple_openhands.plugins.agent_skills.file_reader.file_readers import (
    parse_docx,
    parse_latex,
)


@pytest.fixture(autouse=True)
def reset_current_file() -> None:
    # Reset CURRENT_FILE inside the file_ops module
    import simple_openhands.plugins.agent_skills.file_ops.file_ops as fo

    fo.CURRENT_FILE = None


def _numbered_test_lines(start: int, end: int) -> str:
    return ("\n".join(f"{i}|" for i in range(start, end + 1))) + "\n"


def _generate_test_file_with_lines(temp_path: Path, num_lines: int) -> Path:
    file_path = temp_path / "test_file.py"
    file_path.write_text("\n" * num_lines)
    return file_path


def _calculate_window_bounds(current_line: int, total_lines: int, window_size: int):
    half_window = window_size // 2
    if current_line - half_window < 0:
        start = 1
        end = window_size
    else:
        start = current_line - half_window
        end = current_line + half_window
    return start, end


def _capture_file_operation_error(operation, expected_error_msg: str) -> None:
    with io.StringIO() as buf:
        with contextlib.redirect_stdout(buf):
            operation()
        result = buf.getvalue().strip()
    assert result == expected_error_msg


def test_open_file_unexist_path():
    _capture_file_operation_error(
        lambda: open_file("/unexist/path/a.txt"),
        "ERROR: File /unexist/path/a.txt not found.",
    )


def test_open_file(tmp_path: Path):
    temp_file_path = tmp_path / "a.txt"
    temp_file_path.write_text("Line 1\nLine 2\nLine 3\nLine 4\nLine 5")

    with io.StringIO() as buf:
        with contextlib.redirect_stdout(buf):
            open_file(str(temp_file_path))
        result = buf.getvalue()

    expected = (
        f"[File: {temp_file_path} (5 lines total)]\n"
        "(this is the beginning of the file)\n"
        "1|Line 1\n"
        "2|Line 2\n"
        "3|Line 3\n"
        "4|Line 4\n"
        "5|Line 5\n"
        "(this is the end of the file)\n"
    )
    assert result.split("\n") == expected.split("\n")


def test_open_file_long(tmp_path: Path):
    temp_file_path = tmp_path / "a.txt"
    content = "\n".join([f"Line {i}" for i in range(1, 1001)])
    temp_file_path.write_text(content)

    with io.StringIO() as buf:
        with contextlib.redirect_stdout(buf):
            open_file(str(temp_file_path), 1, 50)
        result = buf.getvalue()

    expected = f"[File: {temp_file_path} (1000 lines total)]\n"
    expected += "(this is the beginning of the file)\n"
    for i in range(1, 51):
        expected += f"{i}|Line {i}\n"
    expected += "(950 more lines below)\n"
    expected += "[Use `scroll_down` to view the next 100 lines of the file!]\n"
    assert result.split("\n") == expected.split("\n")


def test_goto_line(tmp_path: Path):
    temp_file_path = tmp_path / "a.txt"
    total_lines = 1000
    content = "\n".join([f"Line {i}" for i in range(1, total_lines + 1)])
    temp_file_path.write_text(content)

    with io.StringIO() as buf:
        with contextlib.redirect_stdout(buf):
            open_file(str(temp_file_path))
        result = buf.getvalue()

    expected = f"[File: {temp_file_path} ({total_lines} lines total)]\n"
    expected += "(this is the beginning of the file)\n"
    for i in range(1, WINDOW + 1):
        expected += f"{i}|Line {i}\n"
    expected += f"({total_lines - WINDOW} more lines below)\n"
    expected += "[Use `scroll_down` to view the next 100 lines of the file!]\n"
    assert result.split("\n") == expected.split("\n")

    with io.StringIO() as buf:
        with contextlib.redirect_stdout(buf):
            goto_line(500)
        result = buf.getvalue()

    cur_line = 500
    expected = f"[File: {temp_file_path} ({total_lines} lines total)]\n"
    start, end = _calculate_window_bounds(cur_line, total_lines, WINDOW)
    if start == 1:
        expected += "(this is the beginning of the file)\n"
    else:
        expected += f"({start - 1} more lines above)\n"
    for i in range(start, end + 1):
        expected += f"{i}|Line {i}\n"
    if end == total_lines:
        expected += "(this is the end of the file)\n"
    else:
        expected += f"({total_lines - end} more lines below)\n"
    assert result.split("\n") == expected.split("\n")


def test_search_file(tmp_path: Path):
    temp_file_path = tmp_path / "a.txt"
    temp_file_path.write_text("Line 1\nLine 2\nLine 3\nLine 4\nLine 5")

    with io.StringIO() as buf:
        with contextlib.redirect_stdout(buf):
            search_file("Line 5", str(temp_file_path))
        result = buf.getvalue()

    expected = f"[Found 1 matches for \"Line 5\" in {temp_file_path}]\n"
    expected += "Line 5: Line 5\n"
    expected += f"[End of matches for \"Line 5\" in {temp_file_path}]\n"
    assert result.split("\n") == expected.split("\n")


def test_find_file(tmp_path: Path):
    temp_file_path = tmp_path / "a.txt"
    temp_file_path.write_text("Line 1\nLine 2\nLine 3\nLine 4\nLine 5")

    with io.StringIO() as buf:
        with contextlib.redirect_stdout(buf):
            find_file("a.txt", str(tmp_path))
        result = buf.getvalue()

    expected = f"[Found 1 matches for \"a.txt\" in {tmp_path}]\n"
    expected += f"{tmp_path.as_posix()}/a.txt\n"
    expected += f"[End of matches for \"a.txt\" in {tmp_path}]\n"
    assert result.split("\n") == expected.split("\n")


# --------------------- file_reader minimal tests ---------------------


def test_parse_latex(tmp_path: Path):
    test_latex_path = tmp_path / "test.tex"
    with open(test_latex_path, "w", encoding="utf-8") as f:
        f.write(
            r"""
            \documentclass{article}
            \begin{document}
            Hello, this is a test LaTeX document.
            \end{document}
            """
        )

    old_stdout = sys.stdout
    sys.stdout = io.StringIO()
    parse_latex(str(test_latex_path))
    output = sys.stdout.getvalue()
    sys.stdout = old_stdout

    expected_output = (
        f"[Reading LaTex file from {test_latex_path}]\n"
        "Hello, this is a test LaTeX document.\n"
    )
    assert output == expected_output


def test_parse_pptx(tmp_path: Path):
    test_pptx_path = tmp_path / "test.pptx"
    from pptx import Presentation

    pres = Presentation()
    slide1 = pres.slides.add_slide(pres.slide_layouts[0])
    slide1.shapes.title.text = "Hello, this is the first test PPTX slide."
    slide2 = pres.slides.add_slide(pres.slide_layouts[0])
    slide2.shapes.title.text = "Hello, this is the second test PPTX slide."
    pres.save(str(test_pptx_path))

    old_stdout = sys.stdout
    sys.stdout = io.StringIO()
    try:
        from simple_docker_runtime.plugins.agent_skills.file_reader.file_readers import (
            parse_pptx,
        )
    except ModuleNotFoundError:  # container fallback
        from simple_openhands.plugins.agent_skills.file_reader.file_readers import parse_pptx  # type: ignore

    parse_pptx(str(test_pptx_path))
    output = sys.stdout.getvalue()
    sys.stdout = old_stdout

    expected_output = (
        f"[Reading PowerPoint file from {test_pptx_path}]\n"
        "@@ Slide 1 @@\n"
        "Hello, this is the first test PPTX slide.\n\n"
        "@@ Slide 2 @@\n"
        "Hello, this is the second test PPTX slide.\n\n"
    )
    assert output == expected_output


def test_parse_docx(tmp_path: Path):
    import docx

    test_docx_path = tmp_path / "test.docx"
    doc = docx.Document()
    doc.add_paragraph("Hello, this is a test document.")
    doc.add_paragraph("This is the second paragraph.")
    doc.save(str(test_docx_path))

    old_stdout = sys.stdout
    sys.stdout = io.StringIO()
    parse_docx(str(test_docx_path))
    output = sys.stdout.getvalue()
    sys.stdout = old_stdout

    expected_output = (
        f"[Reading DOCX file from {test_docx_path}]\n"
        "@@ Page 1 @@\nHello, this is a test document.\n\n"
        "@@ Page 2 @@\nThis is the second paragraph.\n\n\n"
    )
    assert output == expected_output


def test_parse_pdf(tmp_path: Path):
    pytest.importorskip("reportlab")
    from reportlab.lib.pagesizes import letter  # type: ignore
    from reportlab.pdfgen import canvas  # type: ignore
    try:
        from simple_docker_runtime.plugins.agent_skills.file_reader.file_readers import (
            parse_pdf,
        )
    except ModuleNotFoundError:  # container fallback
        from simple_openhands.plugins.agent_skills.file_reader.file_readers import parse_pdf  # type: ignore

    test_pdf_path = tmp_path / "test.pdf"
    c = canvas.Canvas(str(test_pdf_path), pagesize=letter)
    c.drawString(100, 750, "Hello, this is a test PDF document.")
    c.save()

    old_stdout = sys.stdout
    sys.stdout = io.StringIO()
    parse_pdf(str(test_pdf_path))
    output = sys.stdout.getvalue()
    sys.stdout = old_stdout

    expected_output = (
        f"[Reading PDF file from {test_pdf_path}]\n"
        "@@ Page 1 @@\n"
        "Hello, this is a test PDF document.\n"
    )
    assert output == expected_output


# --------------------- Additional tests from OpenHands ---------------------


def test_open_file_with_indentation(tmp_path: Path):
    temp_file_path = tmp_path / "a.txt"
    temp_file_path.write_text("Line 1\n    Line 2\nLine 3\nLine 4\nLine 5")

    with io.StringIO() as buf:
        with contextlib.redirect_stdout(buf):
            open_file(str(temp_file_path))
        result = buf.getvalue()

    expected = (
        f"[File: {temp_file_path} (5 lines total)]\n"
        "(this is the beginning of the file)\n"
        "1|Line 1\n"
        "2|    Line 2\n"
        "3|Line 3\n"
        "4|Line 4\n"
        "5|Line 5\n"
        "(this is the end of the file)\n"
    )
    assert result.split("\n") == expected.split("\n")


def test_open_file_long_with_lineno(tmp_path: Path):
    temp_file_path = tmp_path / "a.txt"
    content = "\n".join([f"Line {i}" for i in range(1, 1001)])
    temp_file_path.write_text(content)

    cur_line = 100

    with io.StringIO() as buf:
        with contextlib.redirect_stdout(buf):
            open_file(str(temp_file_path), cur_line)
        result = buf.getvalue()

    expected = f"[File: {temp_file_path} (1000 lines total)]\n"
    start, end = _calculate_window_bounds(cur_line, 1000, WINDOW)
    if start == 1:
        expected += "(this is the beginning of the file)\n"
    else:
        expected += f"({start - 1} more lines above)\n"
    for i in range(start, end + 1):
        expected += f"{i}|Line {i}\n"
    if end == 1000:
        expected += "(this is the end of the file)\n"
    else:
        expected += f"({1000 - end} more lines below)\n"
        expected += "[Use `scroll_down` to view the next 100 lines of the file!]\n"
    assert result.split("\n") == expected.split("\n")


def test_goto_line_negative(tmp_path: Path):
    temp_file_path = tmp_path / "a.txt"
    content = "\n".join([f"Line {i}" for i in range(1, 5)])
    temp_file_path.write_text(content)

    with io.StringIO() as buf:
        with contextlib.redirect_stdout(buf):
            open_file(str(temp_file_path))

    _capture_file_operation_error(
        lambda: goto_line(-1), "ERROR: Line number must be between 1 and 4."
    )


def test_goto_line_out_of_bound(tmp_path: Path):
    temp_file_path = tmp_path / "a.txt"
    content = "\n".join([f"Line {i}" for i in range(1, 10)])
    temp_file_path.write_text(content)

    with io.StringIO() as buf:
        with contextlib.redirect_stdout(buf):
            open_file(str(temp_file_path))

    _capture_file_operation_error(
        lambda: goto_line(100), "ERROR: Line number must be between 1 and 9."
    )


def test_scroll_down(tmp_path: Path):
    temp_file_path = tmp_path / "a.txt"
    total_lines = 1000
    content = "\n".join([f"Line {i}" for i in range(1, total_lines + 1)])
    temp_file_path.write_text(content)

    with io.StringIO() as buf:
        with contextlib.redirect_stdout(buf):
            open_file(str(temp_file_path))
        result = buf.getvalue()

    expected = f"[File: {temp_file_path} ({total_lines} lines total)]\n"
    start, end = _calculate_window_bounds(1, total_lines, WINDOW)
    if start == 1:
        expected += "(this is the beginning of the file)\n"
    else:
        expected += f"({start - 1} more lines above)\n"
    for i in range(start, end + 1):
        expected += f"{i}|Line {i}\n"
    if end == total_lines:
        expected += "(this is the end of the file)\n"
    else:
        expected += f"({total_lines - end} more lines below)\n"
        expected += "[Use `scroll_down` to view the next 100 lines of the file!]\n"
    assert result.split("\n") == expected.split("\n")

    with io.StringIO() as buf:
        with contextlib.redirect_stdout(buf):
            scroll_down()
        result = buf.getvalue()

    expected = f"[File: {temp_file_path} ({total_lines} lines total)]\n"
    start = WINDOW + 1
    end = 2 * WINDOW + 1
    if start == 1:
        expected += "(this is the beginning of the file)\n"
    else:
        expected += f"({start - 1} more lines above)\n"
    for i in range(start, end + 1):
        expected += f"{i}|Line {i}\n"
    if end == total_lines:
        expected += "(this is the end of the file)\n"
    else:
        expected += f"({total_lines - end} more lines below)\n"
    assert result.split("\n") == expected.split("\n")


def test_scroll_up(tmp_path: Path):
    temp_file_path = tmp_path / "a.txt"
    total_lines = 1000
    content = "\n".join([f"Line {i}" for i in range(1, total_lines + 1)])
    temp_file_path.write_text(content)

    cur_line = 300
    with io.StringIO() as buf:
        with contextlib.redirect_stdout(buf):
            open_file(str(temp_file_path), cur_line)
        result = buf.getvalue()

    expected = f"[File: {temp_file_path} ({total_lines} lines total)]\n"
    start, end = _calculate_window_bounds(cur_line, total_lines, WINDOW)
    if start == 1:
        expected += "(this is the beginning of the file)\n"
    else:
        expected += f"({start - 1} more lines above)\n"
    for i in range(start, end + 1):
        expected += f"{i}|Line {i}\n"
    if end == total_lines:
        expected += "(this is the end of the file)\n"
    else:
        expected += f"({total_lines - end} more lines below)\n"
        expected += "[Use `scroll_down` to view the next 100 lines of the file!]\n"
    assert result.split("\n") == expected.split("\n")

    with io.StringIO() as buf:
        with contextlib.redirect_stdout(buf):
            scroll_up()
        result = buf.getvalue()

    cur_line = cur_line - WINDOW
    expected = f"[File: {temp_file_path} ({total_lines} lines total)]\n"
    start = cur_line
    end = cur_line + WINDOW

    if start == 1:
        expected += "(this is the beginning of the file)\n"
    else:
        expected += f"({start - 1} more lines above)\n"
    for i in range(start, end + 1):
        expected += f"{i}|Line {i}\n"
    if end == total_lines:
        expected += "(this is the end of the file)\n"
    else:
        expected += f"({total_lines - end} more lines below)\n"
    assert result.split("\n") == expected.split("\n")


def test_scroll_down_edge(tmp_path: Path):
    temp_file_path = tmp_path / "a.txt"
    content = "\n".join([f"Line {i}" for i in range(1, 10)])
    temp_file_path.write_text(content)

    with io.StringIO() as buf:
        with contextlib.redirect_stdout(buf):
            open_file(str(temp_file_path))
        result = buf.getvalue()

    expected = f"[File: {temp_file_path} (9 lines total)]\n"
    expected += "(this is the beginning of the file)\n"
    for i in range(1, 10):
        expected += f"{i}|Line {i}\n"
    expected += "(this is the end of the file)\n"

    with io.StringIO() as buf:
        with contextlib.redirect_stdout(buf):
            scroll_down()
        result = buf.getvalue()

    assert result.split("\n") == expected.split("\n")


def test_print_window_internal(tmp_path: Path):
    test_file_path = tmp_path / "a.txt"
    test_file_path.write_text("")
    open_file(str(test_file_path))
    with open(test_file_path, "w") as file:
        for i in range(1, 101):
            file.write(f"Line `{i}`\n")

    current_line = 50
    window = 2

    with io.StringIO() as buf:
        with contextlib.redirect_stdout(buf):
            _print_window(str(test_file_path), current_line, window, return_str=False)
        result = buf.getvalue()
        expected = (
            "(48 more lines above)\n"
            "49|Line `49`\n"
            "50|Line `50`\n"
            "51|Line `51`\n"
            "(49 more lines below)\n"
        )
        assert result == expected


def test_open_file_large_line_number(tmp_path: Path):
    test_file_path = tmp_path / "a.txt"
    test_file_path.write_text("")
    open_file(str(test_file_path))
    with open(test_file_path, "w") as file:
        for i in range(1, 1000):
            file.write(f"Line `{i}`\n")

    current_line = 800
    window = 100

    with io.StringIO() as buf:
        with contextlib.redirect_stdout(buf):
            open_file(str(test_file_path), current_line, window)
        result = buf.getvalue()
        expected = f"[File: {test_file_path} (999 lines total)]\n"
        expected += "(749 more lines above)\n"
        for i in range(750, 850 + 1):
            expected += f"{i}|Line `{i}`\n"
        expected += "(149 more lines below)\n"
        expected += "[Use `scroll_down` to view the next 100 lines of the file!]\n"
        assert result == expected


def test_search_dir(tmp_path: Path):
    for i in range(1, 101):
        temp_file_path = tmp_path / f"a{i}.txt"
        with open(temp_file_path, "w") as file:
            file.write("Line 1\nLine 2\nLine 3\nLine 4\nLine 5\n")
            if i == 50:
                file.write("bingo")

    with io.StringIO() as buf:
        with contextlib.redirect_stdout(buf):
            search_dir("bingo", str(tmp_path))
        result = buf.getvalue()

    expected = (
        f"[Found 1 matches for \"bingo\" in {tmp_path}]\n"
        f"{tmp_path.as_posix()}/a50.txt (Line 6): bingo\n"
        f"[End of matches for \"bingo\" in {tmp_path}]\n"
    )
    assert result.split("\n") == expected.split("\n")


def test_search_dir_not_exist_term(tmp_path: Path):
    for i in range(1, 101):
        temp_file_path = tmp_path / f"a{i}.txt"
        with open(temp_file_path, "w") as file:
            file.write("Line 1\nLine 2\nLine 3\nLine 4\nLine 5\n")

    with io.StringIO() as buf:
        with contextlib.redirect_stdout(buf):
            search_dir("non-exist", str(tmp_path))
        result = buf.getvalue()

    expected = f"No matches found for \"non-exist\" in {tmp_path}\n"
    assert result.split("\n") == expected.split("\n")


def test_search_dir_too_much_match(tmp_path: Path):
    for i in range(1, 1000):
        temp_file_path = tmp_path / f"a{i}.txt"
        with open(temp_file_path, "w") as file:
            file.write("Line 1\nLine 2\nLine 3\nLine 4\nLine 5\n")

    with io.StringIO() as buf:
        with contextlib.redirect_stdout(buf):
            search_dir("Line 5", str(tmp_path))
        result = buf.getvalue()

    expected = f"More than 999 files matched for \"Line 5\" in {tmp_path}. Please narrow your search.\n"
    assert result.split("\n") == expected.split("\n")


def test_search_dir_cwd(tmp_path: Path, monkeypatch: pytest.MonkeyPatch):
    monkeypatch.chdir(tmp_path)
    for i in range(1, 101):
        temp_file_path = tmp_path / f"a{i}.txt"
        with open(temp_file_path, "w") as file:
            file.write("Line 1\nLine 2\nLine 3\nLine 4\nLine 5\n")
            if i == 50:
                file.write("bingo")

    with io.StringIO() as buf:
        with contextlib.redirect_stdout(buf):
            search_dir("bingo")
        result = buf.getvalue()

    expected = (
        "[Found 1 matches for \"bingo\" in ./]\n"
        "./a50.txt (Line 6): bingo\n"
        "[End of matches for \"bingo\" in ./]\n"
    )
    assert result.split("\n") == expected.split("\n")


def test_search_file_not_exist_term(tmp_path: Path):
    temp_file_path = tmp_path / "a.txt"
    temp_file_path.write_text("Line 1\nLine 2\nLine 3\nLine 4\nLine 5")

    with io.StringIO() as buf:
        with contextlib.redirect_stdout(buf):
            search_file("Line 6", str(temp_file_path))
        result = buf.getvalue()

    expected = f"[No matches found for \"Line 6\" in {temp_file_path}]\n"
    assert result.split("\n") == expected.split("\n")


def test_search_file_not_exist_file():
    _capture_file_operation_error(
        lambda: search_file("Line 6", "/unexist/path/a.txt"),
        "ERROR: File /unexist/path/a.txt not found.",
    )


def test_find_file_cwd(tmp_path: Path, monkeypatch: pytest.MonkeyPatch):
    monkeypatch.chdir(tmp_path)
    temp_file_path = tmp_path / "a.txt"
    temp_file_path.write_text("Line 1\nLine 2\nLine 3\nLine 4\nLine 5")

    with io.StringIO() as buf:
        with contextlib.redirect_stdout(buf):
            find_file("a.txt")
        result = buf.getvalue()
    assert result is not None


def test_find_file_not_exist_file():
    with io.StringIO() as buf:
        with contextlib.redirect_stdout(buf):
            find_file("nonexist.txt")
        result = buf.getvalue()
    expected = '[No matches found for "nonexist.txt" in ./]\n'
    assert result.split('\n') == expected.split('\n')


def test_find_file_not_exist_file_specific_path(tmp_path: Path):
    with io.StringIO() as buf:
        with contextlib.redirect_stdout(buf):
            find_file("nonexist.txt", str(tmp_path))
        result = buf.getvalue()
    expected = f'[No matches found for "nonexist.txt" in {tmp_path}]\n'
    assert result.split('\n') == expected.split('\n')

